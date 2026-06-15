from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from backend.app.schemas.slide import SlideContent
from backend.app.utils.file_utils import generate_temp_path

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ----------------------------
# Helpers
# ----------------------------


def pct(value, total):
    """
    Convert LLM percentage coordinates (0-100)
    into pptx units.
    """
    return total * value / 100


def hex_to_rgb(hex_color):

    if not hex_color:
        return RGBColor(0, 0, 0)

    hex_color = hex_color.replace("#", "")

    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def add_text_element(slide, element, default_color):

    x = pct(element.x, SLIDE_W)
    y = pct(element.y, SLIDE_H)

    w = pct(element.width, SLIDE_W)
    h = pct(element.height, SLIDE_H)

    box = slide.shapes.add_textbox(x, y, w, h)

    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]

    alignment = getattr(element, "alignment", "left")

    if alignment == "center":
        p.alignment = PP_ALIGN.CENTER

    elif alignment == "right":
        p.alignment = PP_ALIGN.RIGHT

    else:
        p.alignment = PP_ALIGN.LEFT

    run = p.add_run()

    run.text = element.content or ""

    size = getattr(element, "font_size", 20)

    if isinstance(size, list):
        size = size[0]

    run.font.size = Pt(size)

    run.font.bold = getattr(element, "bold", False)

    run.font.color.rgb = default_color

    return box


def add_image_placeholder(slide, element):

    x = pct(element.x, SLIDE_W)
    y = pct(element.y, SLIDE_H)

    w = pct(element.width, SLIDE_W)
    h = pct(element.height, SLIDE_H)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)

    shape.text = "IMAGE\n\n" + str(element.content)

    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


# ----------------------------
# Background
# ----------------------------


def apply_background(slide, slide_data):

    fill = slide.background.fill
    fill.solid()

    color = "#FFFFFF"

    if hasattr(slide_data, "background"):

        if slide_data.background:
            color = getattr(slide_data.background, "color", "#FFFFFF")

    fill.fore_color.rgb = hex_to_rgb(color)


# ----------------------------
# Main Builder
# ----------------------------


def build_pptx(
    slides: List[SlideContent], presentation_title: str, theme_name="default"
):

    prs = Presentation()

    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    for s in slides:

        slide = prs.slides.add_slide(blank)

        apply_background(slide, s)

        #
        # Theme colors
        #

        title_color = RGBColor(17, 17, 17)

        body_color = RGBColor(50, 50, 50)

        #
        # NEW SYSTEM:
        # Render LLM elements directly
        #

        elements = getattr(s, "elements", [])

        if elements:

            for element in elements:

                if element.type == "text":

                    add_text_element(slide, element, body_color)

                elif element.type == "image":

                    add_image_placeholder(slide, element)

        #
        # FALLBACK:
        # Old bullet schema
        #

        else:

            add_text_element(
                slide,
                type(
                    "obj",
                    (),
                    {
                        "x": 5,
                        "y": 5,
                        "width": 90,
                        "height": 15,
                        "content": s.title,
                        "font_size": getattr(s, "title_size", 32),
                        "bold": True,
                        "alignment": "center",
                    },
                )(),
                title_color,
            )

            if s.bullets:

                bullets = "\n".join("• " + b for b in s.bullets)

                add_text_element(
                    slide,
                    type(
                        "obj",
                        (),
                        {
                            "x": 10,
                            "y": 30,
                            "width": 80,
                            "height": 50,
                            "content": bullets,
                            "font_size": getattr(s, "body_size", 20),
                            "bold": False,
                            "alignment": "left",
                        },
                    )(),
                    body_color,
                )

        #
        # Speaker notes
        #

        if getattr(s, "speaker_notes", None):

            slide.notes_slide.notes_text_frame.text = s.speaker_notes

    path = generate_temp_path(".pptx")

    prs.save(path)

    return path
